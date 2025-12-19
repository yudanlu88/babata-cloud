import streamlit as st
from openai import OpenAI
from docx import Document
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import io
import asyncio
import edge_tts
import os
from datetime import datetime
from pptx import Presentation # 🔥 新增：PPT 核心库
from pptx.util import Inches, Pt

# --- 1. UI 深度定制 (C路线核心) ---
st.set_page_config(page_title="巴巴塔·灵感中枢", page_icon="🪐", layout="wide")

# 注入极客风格 CSS (黑金/霓虹风)
st.markdown("""
<style>
    /* 全局背景微调 */
    .stApp { background-color: #f0f2f6; }
    
    /* 侧边栏美化 */
    [data-testid="stSidebar"] {
        background-color: #111827; /* 深色侧边栏 */
        color: white;
    }
    [data-testid="stSidebar"] h1, [data-testid="stSidebar"] p, [data-testid="stSidebar"] label {
        color: #e0e0e0 !important;
    }
    
    /* 按钮高级特效 (渐变+阴影) */
    .stButton>button {
        background: linear-gradient(45deg, #2563eb, #9333ea);
        color: white;
        border: none;
        border-radius: 8px;
        height: 50px;
        font-weight: bold;
        font-size: 16px;
        box-shadow: 0 4px 15px rgba(37, 99, 235, 0.4);
        transition: all 0.3s ease;
    }
    .stButton>button:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(147, 51, 234, 0.6);
    }
    
    /* 标题字体优化 */
    h1 { font-family: 'Helvetica Neue', sans-serif; font-weight: 800; color: #1f2937; }
</style>
""", unsafe_allow_html=True)

# --- 2. 配置检查 ---
if "DEEPSEEK_KEY" in st.secrets:
    api_key = st.secrets["DEEPSEEK_KEY"]
else:
    st.error("⚠️ 警告：请先在 Streamlit 后台配置 Secrets！")
    st.stop()

client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com")

# --- 3. 核心功能函数 ---

# (A) 数据库模块 (保留之前的记忆功能)
DB_FILE = "babata_memory.csv"
def save_to_db(mode, topic, content):
    if not os.path.exists(DB_FILE):
        pd.DataFrame(columns=["时间", "模式", "主题", "摘要"]).to_csv(DB_FILE, index=False)
    summary = content[:30].replace("#", "").replace("*", "") + "..."
    new_data = pd.DataFrame([{
        "时间": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "模式": mode, "主题": topic, "摘要": summary
    }])
    new_data.to_csv(DB_FILE, mode='a', header=False, index=False)

def load_from_db():
    try: return pd.read_csv(DB_FILE)
    except: return pd.DataFrame()

# (B) 语音模块
async def generate_audio_file(text, filename="output.mp3"):
    communicate = edge_tts.Communicate(text, "zh-CN-XiaoxiaoNeural")
    await communicate.save(filename)

# (C) PPT 生成引擎 🔥 (B路线核心)
def create_ppt(topic, full_text):
    prs = Presentation()
    
    # 1. 封面页
    slide_layout = prs.slide_layouts[0] # 标题页布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d')}\nBy 巴巴塔万能助手"
    
    # 2. 内容页 (智能切分)
    # 我们假设 AI 输出是按 Markdown 标题分段的，或者直接把整段放进去
    # 这里做一个简单的处理：每 500 字一页，或者按段落
    
    slide_layout = prs.slide_layouts[1] # 标题+内容布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "核心方案详情"
    
    # 清洗 Markdown 符号，让 PPT 看起来干净点
    clean_text = full_text.replace("### ", "").replace("**", "").replace("## ", "")
    
    # 这里为了防止文字溢出，只取前1000字，或者你可以做更复杂的切页逻辑
    tf = body.text_frame
    tf.text = clean_text[:800] + "..." # 简单截断，防止撑爆PPT
    
    # 保存到内存
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- 4. 侧边栏 ---
with st.sidebar:
    st.image("https://api.iconify.design/vscode-icons:file-type-python.svg", width=50)
    st.title("控制台")
    
    app_mode = st.selectbox("功能模式", 
        ["💼 商业策划案", "📕 小红书爆款", "📊 职场周报大师", "❤️ 情感/哄人专家"]
    )
    st.divider()
    
    # 商业模式才显示行业
    if app_mode == "💼 商业策划案":
        industry = st.selectbox("行业赛道", ["🚀 AI/科技", "🛒 消费/零售", "🏥 医疗", "⚙️ 制造"])
    
    style_mode = st.radio("AI 语气", ["专业理性", "毒舌巴巴塔", "温柔贴心", "热情激昂"])
    word_count = st.slider("字数", 200, 3000, 800)
    enable_voice = st.toggle("🔊 语音播报", value=True)

    # 历史记录入口
    st.divider()
    with st.expander("📜 历史档案"):
        df = load_from_db()
        if not df.empty:
            st.dataframe(df[["时间", "主题"]], hide_index=True)
        else:
            st.caption("暂无记录")

# --- 5. 主逻辑 ---
def get_prompt(mode):
    if mode == "💼 商业策划案": return "【强制中文】输出商业策划案。结构：🎯摘要、⚡痛点、💎方案、💰模式。"
    elif mode == "📕 小红书爆款": return "小红书爆款博主。emoji多，标题抓眼球。"
    elif mode == "📊 职场周报大师": return "互联网大厂P8。多用黑话：赋能、闭环。"
    else: return "情感专家。温柔体贴。"

st.title(f"{app_mode} 🪐")

with st.form("main_form"):
    if app_mode == "💼 商业策划案": ph = "输入项目点子..."
    else: ph = "输入主题..."
    user_input = st.text_input("💡 灵感输入", placeholder=ph)
    # 按钮会自动应用上面的 CSS 特效
    submitted = st.form_submit_button("🚀 启动引擎")

if submitted and user_input:
    out = st.empty()
    full_text = ""
    prompt_sys = get_prompt(app_mode)
    
    try:
        # 1. AI 生成
        stream = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": f"{prompt_sys}\n语气:{style_mode}\n字数:{word_count}"},
                {"role": "user", "content": user_input}
            ],
            stream=True
        )
        for chunk in stream:
            if chunk.choices[0].delta.content:
                full_text += chunk.choices[0].delta.content
                out.markdown(full_text + "▌")
        out.markdown(full_text)
        
        # 2. 存入数据库
        save_to_db(app_mode, user_input, full_text)
        
        # 3. 语音 & 文件处理
        if enable_voice:
            asyncio.run(generate_audio_file(full_text.replace("#",""), "voice.mp3"))
            st.audio("voice.mp3", autoplay=True)
            
        # 🔥 4. 生成 PPT (内存流)
        ppt_file = create_ppt(user_input, full_text)
        
        # 生成 Word (内存流)
        doc = Document()
        doc.add_heading(user_input, 0)
        doc.add_paragraph(full_text)
        doc_io = io.BytesIO()
        doc.save(doc_io)
        doc_io.seek(0)
        
        # --- 下载区 ---
        st.divider()
        st.success("✅ 任务完成！已生成全套资料")
        
        c1, c2, c3 = st.columns(3)
        # 漂亮的下载按钮
        c1.download_button("📘 下载 Word", doc_io, f"{user_input}.docx", use_container_width=True)
        # 🔥 新增 PPT 下载按钮
        c2.download_button("📊 下载 PPT", ppt_file, f"{user_input}.pptx", use_container_width=True)
        # Markdown
        c3.download_button("📝 下载 MD", full_text, f"{user_input}.md", use_container_width=True)

    except Exception as e:
        st.error(f"发生错误: {e}")
